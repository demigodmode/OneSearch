# Copyright (C) 2025 demigodmode
# SPDX-License-Identifier: AGPL-3.0-only

"""
Office document extractors for Microsoft Office formats
Supports .docx, .xlsx, and .pptx files
"""
import logging
from pathlib import Path
from typing import Optional

from .base import BaseExtractor, extractor_registry
from ..schemas import Document
from ..config import settings

logger = logging.getLogger(__name__)

# Keywords that indicate password-protected or encrypted files
# across different library versions and error messages
PASSWORD_KEYWORDS = frozenset([
    "password",
    "encrypted",
    "protected",
    "decrypt",
    "cipher",
    "encryption",
])


def _is_password_protected_error(error: Exception) -> bool:
    """
    Check if an exception indicates a password-protected or encrypted file.

    Uses keyword matching against the error message and type name to detect
    password protection across different library versions.

    Args:
        error: The exception to check

    Returns:
        True if the error indicates password protection, False otherwise
    """
    error_str = str(error).lower()
    error_type = type(error).__name__.lower()

    # Check error message and type name for password-related keywords
    for keyword in PASSWORD_KEYWORDS:
        if keyword in error_str or keyword in error_type:
            return True

    return False


class DocxExtractor(BaseExtractor):
    """
    Extractor for Microsoft Word documents (.docx)

    Extracts:
    - All paragraph text
    - Table content
    - Metadata (title, author, word count)
    """

    SUPPORTED_EXTENSIONS = ['.docx']

    MAX_FILE_SIZE = settings.max_office_file_size_mb * 1024 * 1024
    TIMEOUT = settings.office_extraction_timeout

    def extract(self, file_path: str) -> Document:
        """
        Extract text and metadata from DOCX file

        Args:
            file_path: Path to DOCX file

        Returns:
            Document with extracted content

        Raises:
            FileNotFoundError: If file doesn't exist
            ValueError: If file is too large or password-protected
        """
        from docx import Document as DocxDocument
        from docx.opc.exceptions import PackageNotFoundError

        try:
            self._check_file_size(file_path)

            try:
                content, metadata = self._extract_docx_content(file_path)
            except PackageNotFoundError as e:
                logger.warning(f"DOCX file corrupted or invalid: {file_path}: {e}")
                content = ""
                metadata = {
                    "extraction_error": "File corrupted or invalid format",
                    "extraction_failed": True,
                }
            except Exception as e:
                if _is_password_protected_error(e):
                    logger.warning(f"DOCX file is password-protected: {file_path}")
                    content = ""
                    metadata = {
                        "extraction_error": "File is password-protected",
                        "extraction_failed": True,
                    }
                else:
                    logger.warning(f"DOCX extraction failed for {file_path}: {e}")
                    content = ""
                    metadata = {
                        "extraction_error": str(e),
                        "extraction_failed": True,
                    }

            doc = self._create_base_document(file_path, content)
            doc.title = metadata.get("title") or Path(file_path).stem
            doc.metadata = metadata

            return doc

        except (FileNotFoundError, ValueError):
            raise
        except Exception as e:
            logger.error(f"Unexpected error extracting DOCX from {file_path}: {e}", exc_info=True)
            raise

    def _extract_docx_content(self, file_path: str) -> tuple[str, dict]:
        """Extract text content and metadata from DOCX"""
        from docx import Document as DocxDocument

        metadata = {}
        text_parts = []

        doc = DocxDocument(file_path)

        # Extract core properties
        core_props = doc.core_properties
        if core_props.title:
            metadata["title"] = core_props.title
        if core_props.author:
            metadata["author"] = core_props.author
        if core_props.subject:
            metadata["subject"] = core_props.subject
        if core_props.keywords:
            metadata["keywords"] = core_props.keywords

        # Extract paragraph text
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)

        # Extract table content
        table_count = 0
        for table in doc.tables:
            table_count += 1
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text_parts.append(" | ".join(row_text))

        full_text = "\n".join(text_parts)

        # Add extraction stats
        metadata["paragraph_count"] = len(doc.paragraphs)
        metadata["table_count"] = table_count
        metadata["extracted_text_length"] = len(full_text)
        metadata["extraction_failed"] = False

        if not full_text.strip():
            metadata["extraction_warning"] = "No text extracted (document may be empty or image-based)"

        return full_text, metadata

    def _get_document_type(self) -> str:
        return "docx"


class XlsxExtractor(BaseExtractor):
    """
    Extractor for Microsoft Excel spreadsheets (.xlsx)

    Extracts:
    - All cell values as text
    - Sheet names and structure metadata
    """

    SUPPORTED_EXTENSIONS = ['.xlsx']

    MAX_FILE_SIZE = settings.max_office_file_size_mb * 1024 * 1024
    TIMEOUT = settings.office_extraction_timeout

    # Max rows/columns to extract to prevent memory issues
    MAX_ROWS = 10000
    MAX_COLS = 100

    def extract(self, file_path: str) -> Document:
        """
        Extract text and metadata from XLSX file

        Args:
            file_path: Path to XLSX file

        Returns:
            Document with extracted content
        """
        from openpyxl import load_workbook
        from openpyxl.utils.exceptions import InvalidFileException

        try:
            self._check_file_size(file_path)

            try:
                content, metadata = self._extract_xlsx_content(file_path)
            except InvalidFileException as e:
                logger.warning(f"XLSX file corrupted or invalid: {file_path}: {e}")
                content = ""
                metadata = {
                    "extraction_error": "File corrupted or invalid format",
                    "extraction_failed": True,
                }
            except Exception as e:
                if _is_password_protected_error(e):
                    logger.warning(f"XLSX file is password-protected: {file_path}")
                    content = ""
                    metadata = {
                        "extraction_error": "File is password-protected",
                        "extraction_failed": True,
                    }
                else:
                    logger.warning(f"XLSX extraction failed for {file_path}: {e}")
                    content = ""
                    metadata = {
                        "extraction_error": str(e),
                        "extraction_failed": True,
                    }

            doc = self._create_base_document(file_path, content)
            doc.title = metadata.get("title") or Path(file_path).stem
            doc.metadata = metadata

            return doc

        except (FileNotFoundError, ValueError):
            raise
        except Exception as e:
            logger.error(f"Unexpected error extracting XLSX from {file_path}: {e}", exc_info=True)
            raise

    def _extract_xlsx_content(self, file_path: str) -> tuple[str, dict]:
        """Extract text content and metadata from XLSX"""
        from openpyxl import load_workbook

        metadata = {}
        text_parts = []

        # Load workbook in read-only mode for better memory usage
        wb = load_workbook(file_path, read_only=True, data_only=True)

        try:
            # Get sheet names
            sheet_names = wb.sheetnames
            metadata["sheet_names"] = sheet_names
            metadata["sheet_count"] = len(sheet_names)

            total_rows = 0
            total_cells = 0

            for sheet_name in sheet_names:
                sheet = wb[sheet_name]
                text_parts.append(f"=== Sheet: {sheet_name} ===")

                row_count = 0
                for row in sheet.iter_rows(max_row=self.MAX_ROWS, max_col=self.MAX_COLS):
                    row_values = []
                    for cell in row:
                        if cell.value is not None:
                            # Convert to string, handling various types
                            cell_str = str(cell.value).strip()
                            if cell_str:
                                row_values.append(cell_str)
                                total_cells += 1

                    if row_values:
                        text_parts.append(" | ".join(row_values))
                        row_count += 1

                    if row_count >= self.MAX_ROWS:
                        text_parts.append(f"... (truncated at {self.MAX_ROWS} rows)")
                        break

                total_rows += row_count

            full_text = "\n".join(text_parts)
        finally:
            wb.close()

        metadata["total_rows_extracted"] = total_rows
        metadata["total_cells_extracted"] = total_cells
        metadata["extracted_text_length"] = len(full_text)
        metadata["extraction_failed"] = False

        if total_cells == 0:
            metadata["extraction_warning"] = "No cell values extracted (spreadsheet may be empty)"

        return full_text, metadata

    def _get_document_type(self) -> str:
        return "xlsx"


class PptxExtractor(BaseExtractor):
    """
    Extractor for Microsoft PowerPoint presentations (.pptx)

    Extracts:
    - All slide text content
    - Speaker notes
    - Metadata (title, author, slide count)
    """

    SUPPORTED_EXTENSIONS = ['.pptx']

    MAX_FILE_SIZE = settings.max_office_file_size_mb * 1024 * 1024
    TIMEOUT = settings.office_extraction_timeout

    def extract(self, file_path: str) -> Document:
        """
        Extract text and metadata from PPTX file

        Args:
            file_path: Path to PPTX file

        Returns:
            Document with extracted content
        """
        from pptx import Presentation
        from pptx.exc import PackageNotFoundError

        try:
            self._check_file_size(file_path)

            try:
                content, metadata = self._extract_pptx_content(file_path)
            except PackageNotFoundError as e:
                logger.warning(f"PPTX file corrupted or invalid: {file_path}: {e}")
                content = ""
                metadata = {
                    "extraction_error": "File corrupted or invalid format",
                    "extraction_failed": True,
                }
            except Exception as e:
                if _is_password_protected_error(e):
                    logger.warning(f"PPTX file is password-protected: {file_path}")
                    content = ""
                    metadata = {
                        "extraction_error": "File is password-protected",
                        "extraction_failed": True,
                    }
                else:
                    logger.warning(f"PPTX extraction failed for {file_path}: {e}")
                    content = ""
                    metadata = {
                        "extraction_error": str(e),
                        "extraction_failed": True,
                    }

            doc = self._create_base_document(file_path, content)
            doc.title = metadata.get("title") or Path(file_path).stem
            doc.metadata = metadata

            return doc

        except (FileNotFoundError, ValueError):
            raise
        except Exception as e:
            logger.error(f"Unexpected error extracting PPTX from {file_path}: {e}", exc_info=True)
            raise

    def _extract_pptx_content(self, file_path: str) -> tuple[str, dict]:
        """Extract text content and metadata from PPTX"""
        from pptx import Presentation

        metadata = {}
        text_parts = []

        prs = Presentation(file_path)

        # Extract core properties
        core_props = prs.core_properties
        if core_props.title:
            metadata["title"] = core_props.title
        if core_props.author:
            metadata["author"] = core_props.author
        if core_props.subject:
            metadata["subject"] = core_props.subject
        if core_props.keywords:
            metadata["keywords"] = core_props.keywords

        # Get slide count
        slide_count = len(prs.slides)
        metadata["slide_count"] = slide_count

        # Extract text from each slide
        slides_with_text = 0
        notes_count = 0

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text_parts = []

            # Extract text from shapes
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        para_text = paragraph.text.strip()
                        if para_text:
                            slide_text_parts.append(para_text)

                # Extract text from tables in shapes
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            slide_text_parts.append(" | ".join(row_text))

            if slide_text_parts:
                text_parts.append(f"=== Slide {slide_num} ===")
                text_parts.extend(slide_text_parts)
                slides_with_text += 1

            # Extract speaker notes
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame and notes_frame.text.strip():
                    text_parts.append(f"[Notes for Slide {slide_num}]")
                    text_parts.append(notes_frame.text.strip())
                    notes_count += 1

        full_text = "\n".join(text_parts)

        metadata["slides_with_text"] = slides_with_text
        metadata["slides_with_notes"] = notes_count
        metadata["extracted_text_length"] = len(full_text)
        metadata["extraction_failed"] = False

        if not full_text.strip():
            metadata["extraction_warning"] = "No text extracted (presentation may be empty or image-based)"

        return full_text, metadata

    def _get_document_type(self) -> str:
        return "pptx"


# Register all Office extractors
extractor_registry.register(DocxExtractor)
extractor_registry.register(XlsxExtractor)
extractor_registry.register(PptxExtractor)
