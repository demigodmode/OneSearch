# Copyright (C) 2025 demigodmode
# SPDX-License-Identifier: AGPL-3.0-only

"""
Tests for document extractors
"""
import pytest
import asyncio
from pathlib import Path
from tempfile import TemporaryDirectory
from pypdf import PdfWriter

from app.extractors import (
    TextExtractor,
    MarkdownExtractor,
    PDFExtractor,
    DocxExtractor,
    XlsxExtractor,
    PptxExtractor,
    extractor_registry,
)


@pytest.fixture
def temp_dir():
    """Create temporary directory for test files"""
    with TemporaryDirectory() as tmp:
        yield Path(tmp)


@pytest.fixture
def sample_text_file(temp_dir):
    """Create sample text file"""
    file_path = temp_dir / "sample.txt"
    file_path.write_text("This is a test file.\nIt has multiple lines.\nFor testing purposes.")
    return str(file_path)


@pytest.fixture
def sample_markdown_file(temp_dir):
    """Create sample markdown file with front-matter"""
    file_path = temp_dir / "sample.md"
    content = """---
title: Test Document
tags: [test, markdown]
author: Test Author
---

# Main Heading

This is the content of the markdown file.

## Subheading

More content here.
"""
    file_path.write_text(content)
    return str(file_path)


@pytest.fixture
def sample_pdf_file(temp_dir):
    """Create sample PDF file"""
    file_path = temp_dir / "sample.pdf"

    # Create a simple PDF with PyPDF2
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)

    # Add metadata
    writer.add_metadata({
        "/Title": "Test PDF Document",
        "/Author": "Test Author",
    })

    with open(file_path, "wb") as f:
        writer.write(f)

    return str(file_path)


class TestTextExtractor:
    """Tests for TextExtractor"""

    @pytest.mark.asyncio
    async def test_extract_text_file(self, sample_text_file):
        """Test extracting from a plain text file"""
        extractor = TextExtractor("test_source", "Test Source")
        doc = await extractor.extract_with_timeout(sample_text_file)

        assert doc.source_id == "test_source"
        assert doc.source_name == "Test Source"
        assert doc.type == "text"
        assert doc.basename == "sample.txt"
        assert doc.extension == "txt"
        assert "This is a test file" in doc.content
        assert doc.metadata["line_count"] == 3
        assert doc.title is not None

    @pytest.mark.asyncio
    async def test_supports_file(self, temp_dir):
        """Test file extension support"""
        assert TextExtractor.supports_file(str(temp_dir / "test.txt"))
        assert TextExtractor.supports_file(str(temp_dir / "test.log"))
        assert TextExtractor.supports_file(str(temp_dir / "test.py"))
        assert not TextExtractor.supports_file(str(temp_dir / "test.pdf"))

    @pytest.mark.asyncio
    async def test_file_size_limit(self, temp_dir):
        """Test file size limit enforcement"""
        # Create large file
        large_file = temp_dir / "large.txt"
        large_file.write_text("x" * (11 * 1024 * 1024))  # 11 MB

        extractor = TextExtractor("test_source", "Test Source")

        with pytest.raises(ValueError, match="File too large"):
            await extractor.extract_with_timeout(str(large_file))

    @pytest.mark.asyncio
    async def test_encoding_detection(self, temp_dir):
        """Test encoding detection for non-UTF8 files"""
        # Create file with latin-1 encoding
        latin1_file = temp_dir / "latin1.txt"
        with open(latin1_file, "w", encoding="latin-1") as f:
            f.write("Test with special chars: \xe9\xe8")

        extractor = TextExtractor("test_source", "Test Source")
        doc = await extractor.extract_with_timeout(str(latin1_file))

        assert doc.content is not None
        assert len(doc.content) > 0

    @pytest.mark.asyncio
    async def test_timeout(self, temp_dir, monkeypatch):
        """Test extraction timeout"""
        import time

        # Create a mock extract that does blocking I/O for too long
        def slow_extract(self, file_path):
            time.sleep(10)  # Blocking sleep
            return None

        monkeypatch.setattr(TextExtractor, "extract", slow_extract)

        extractor = TextExtractor("test_source", "Test Source")
        extractor.TIMEOUT = 1  # 1 second timeout

        test_file = temp_dir / "test.txt"
        test_file.write_text("test")

        with pytest.raises(TimeoutError):
            await extractor.extract_with_timeout(str(test_file))


class TestMarkdownExtractor:
    """Tests for MarkdownExtractor"""

    @pytest.mark.asyncio
    async def test_extract_markdown_with_frontmatter(self, sample_markdown_file):
        """Test extracting markdown with YAML front-matter"""
        extractor = MarkdownExtractor("test_source", "Test Source")
        doc = await extractor.extract_with_timeout(sample_markdown_file)

        assert doc.source_id == "test_source"
        assert doc.type == "markdown"
        assert doc.title == "Test Document"
        assert "Main Heading" in doc.content
        assert doc.metadata["has_frontmatter"] is True
        assert doc.metadata["tags"] == ["test", "markdown"]
        assert doc.metadata["author"] == "Test Author"

    @pytest.mark.asyncio
    async def test_extract_markdown_without_frontmatter(self, temp_dir):
        """Test extracting markdown without front-matter"""
        file_path = temp_dir / "no_frontmatter.md"
        file_path.write_text("# Title from Heading\n\nContent here.")

        extractor = MarkdownExtractor("test_source", "Test Source")
        doc = await extractor.extract_with_timeout(str(file_path))

        assert doc.title == "Title from Heading"
        assert doc.metadata["has_frontmatter"] is False

    @pytest.mark.asyncio
    async def test_markdown_fallback_title(self, temp_dir):
        """Test fallback to filename for title"""
        file_path = temp_dir / "my_document.md"
        file_path.write_text("Just some content without a heading.")

        extractor = MarkdownExtractor("test_source", "Test Source")
        doc = await extractor.extract_with_timeout(str(file_path))

        assert doc.title == "my_document"

    @pytest.mark.asyncio
    async def test_supports_file(self, temp_dir):
        """Test file extension support"""
        assert MarkdownExtractor.supports_file(str(temp_dir / "test.md"))
        assert MarkdownExtractor.supports_file(str(temp_dir / "test.markdown"))
        assert not MarkdownExtractor.supports_file(str(temp_dir / "test.txt"))


class TestPDFExtractor:
    """Tests for PDFExtractor"""

    @pytest.mark.asyncio
    async def test_extract_pdf(self, sample_pdf_file):
        """Test extracting from PDF file"""
        extractor = PDFExtractor("test_source", "Test Source")
        doc = await extractor.extract_with_timeout(sample_pdf_file)

        assert doc.source_id == "test_source"
        assert doc.type == "pdf"
        assert doc.basename == "sample.pdf"
        assert doc.extension == "pdf"
        assert doc.metadata["page_count"] == 1
        assert doc.metadata.get("title") == "Test PDF Document"

    @pytest.mark.asyncio
    async def test_supports_file(self, temp_dir):
        """Test file extension support"""
        assert PDFExtractor.supports_file(str(temp_dir / "test.pdf"))
        assert not PDFExtractor.supports_file(str(temp_dir / "test.txt"))

    @pytest.mark.asyncio
    async def test_corrupted_pdf_fallback(self, temp_dir):
        """Test fallback behavior on corrupted PDF"""
        # Create invalid PDF file
        corrupted_pdf = temp_dir / "corrupted.pdf"
        corrupted_pdf.write_text("This is not a PDF file")

        extractor = PDFExtractor("test_source", "Test Source")
        doc = await extractor.extract_with_timeout(str(corrupted_pdf))

        # Should still create document with minimal content
        assert doc is not None
        assert doc.content == ""
        assert doc.metadata.get("extraction_failed") is True
        assert "extraction_error" in doc.metadata

    @pytest.mark.asyncio
    async def test_encrypted_pdf(self, temp_dir):
        """Test handling of encrypted PDF that cannot be decrypted"""
        # Create encrypted PDF
        from pypdf import PdfWriter, PdfReader

        encrypted_pdf = temp_dir / "encrypted.pdf"

        # Create and encrypt a PDF
        writer = PdfWriter()
        writer.add_blank_page(width=200, height=200)
        writer.encrypt(user_password="password123", owner_password="password456")

        with open(encrypted_pdf, "wb") as f:
            writer.write(f)

        extractor = PDFExtractor("test_source", "Test Source")
        doc = await extractor.extract_with_timeout(str(encrypted_pdf))

        # Should handle gracefully with error in metadata
        assert doc is not None
        assert doc.metadata.get("extraction_failed") is True
        error_msg = doc.metadata.get("extraction_error", "").lower()
        # pypdf returns "file has not been decrypted" for encrypted PDFs
        assert "decrypt" in error_msg or "encrypted" in error_msg


class TestExtractorRegistry:
    """Tests for ExtractorRegistry"""

    def test_get_extractor_for_text(self, temp_dir):
        """Test getting extractor for text file"""
        file_path = temp_dir / "test.txt"
        extractor = extractor_registry.get_extractor(
            str(file_path), "source1", "Source 1"
        )

        assert extractor is not None
        assert isinstance(extractor, TextExtractor)

    def test_get_extractor_for_markdown(self, temp_dir):
        """Test getting extractor for markdown file"""
        file_path = temp_dir / "test.md"
        extractor = extractor_registry.get_extractor(
            str(file_path), "source1", "Source 1"
        )

        assert extractor is not None
        assert isinstance(extractor, MarkdownExtractor)

    def test_get_extractor_for_pdf(self, temp_dir):
        """Test getting extractor for PDF file"""
        file_path = temp_dir / "test.pdf"
        extractor = extractor_registry.get_extractor(
            str(file_path), "source1", "Source 1"
        )

        assert extractor is not None
        assert isinstance(extractor, PDFExtractor)

    def test_get_extractor_unsupported(self, temp_dir):
        """Test getting extractor for unsupported file type"""
        file_path = temp_dir / "test.xyz"
        extractor = extractor_registry.get_extractor(
            str(file_path), "source1", "Source 1"
        )

        assert extractor is None

    def test_get_supported_extensions(self):
        """Test getting list of supported extensions"""
        extensions = extractor_registry.get_supported_extensions()

        assert ".txt" in extensions
        assert ".md" in extensions
        assert ".pdf" in extensions
        assert ".docx" in extensions
        assert ".xlsx" in extensions
        assert ".pptx" in extensions
        assert len(extensions) > 10  # Should have many supported types


# Fixtures for Office files
@pytest.fixture
def sample_docx_file(temp_dir):
    """Create sample DOCX file"""
    from docx import Document as DocxDocument

    file_path = temp_dir / "sample.docx"
    doc = DocxDocument()

    # Set core properties
    doc.core_properties.title = "Test Word Document"
    doc.core_properties.author = "Test Author"

    # Add content
    doc.add_heading("Main Heading", level=1)
    doc.add_paragraph("This is the first paragraph of the document.")
    doc.add_paragraph("This is the second paragraph with more content.")

    # Add a table
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Cell A1"
    table.cell(0, 1).text = "Cell B1"
    table.cell(1, 0).text = "Cell A2"
    table.cell(1, 1).text = "Cell B2"

    doc.save(file_path)
    return str(file_path)


@pytest.fixture
def sample_xlsx_file(temp_dir):
    """Create sample XLSX file"""
    from openpyxl import Workbook

    file_path = temp_dir / "sample.xlsx"
    wb = Workbook()

    # First sheet
    ws1 = wb.active
    ws1.title = "Data Sheet"
    ws1["A1"] = "Name"
    ws1["B1"] = "Value"
    ws1["A2"] = "Item One"
    ws1["B2"] = 100
    ws1["A3"] = "Item Two"
    ws1["B3"] = 200

    # Second sheet
    ws2 = wb.create_sheet("Summary")
    ws2["A1"] = "Total Items"
    ws2["B1"] = 2

    wb.save(file_path)
    return str(file_path)


@pytest.fixture
def sample_pptx_file(temp_dir):
    """Create sample PPTX file"""
    from pptx import Presentation
    from pptx.util import Inches

    file_path = temp_dir / "sample.pptx"
    prs = Presentation()

    # Set core properties
    prs.core_properties.title = "Test Presentation"
    prs.core_properties.author = "Test Author"

    # Add title slide
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    title.text = "Presentation Title"
    subtitle.text = "Subtitle text here"

    # Add content slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide2.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Slide Two Title"
    tf = body_shape.text_frame
    tf.text = "First bullet point"
    p = tf.add_paragraph()
    p.text = "Second bullet point"

    prs.save(file_path)
    return str(file_path)


class TestDocxExtractor:
    """Tests for DocxExtractor"""

    @pytest.mark.asyncio
    async def test_extract_docx(self, sample_docx_file):
        """Test extracting from DOCX file"""
        extractor = DocxExtractor("test_source", "Test Source")
        doc = await extractor.extract_with_timeout(sample_docx_file)

        assert doc.source_id == "test_source"
        assert doc.type == "docx"
        assert doc.basename == "sample.docx"
        assert doc.extension == "docx"
        assert doc.title == "Test Word Document"
        assert "first paragraph" in doc.content
        assert "second paragraph" in doc.content
        assert "Cell A1" in doc.content
        assert doc.metadata.get("author") == "Test Author"
        assert doc.metadata.get("extraction_failed") is False

    @pytest.mark.asyncio
    async def test_supports_file(self, temp_dir):
        """Test file extension support"""
        assert DocxExtractor.supports_file(str(temp_dir / "test.docx"))
        assert not DocxExtractor.supports_file(str(temp_dir / "test.doc"))
        assert not DocxExtractor.supports_file(str(temp_dir / "test.txt"))

    @pytest.mark.asyncio
    async def test_corrupted_docx_fallback(self, temp_dir):
        """Test fallback behavior on corrupted DOCX"""
        corrupted_docx = temp_dir / "corrupted.docx"
        corrupted_docx.write_text("This is not a DOCX file")

        extractor = DocxExtractor("test_source", "Test Source")
        doc = await extractor.extract_with_timeout(str(corrupted_docx))

        assert doc is not None
        assert doc.content == ""
        assert doc.metadata.get("extraction_failed") is True
        assert "extraction_error" in doc.metadata

    @pytest.mark.asyncio
    async def test_docx_without_title(self, temp_dir):
        """Test DOCX without title falls back to filename"""
        from docx import Document as DocxDocument

        file_path = temp_dir / "no_title.docx"
        doc = DocxDocument()
        doc.add_paragraph("Just some content")
        doc.save(file_path)

        extractor = DocxExtractor("test_source", "Test Source")
        result = await extractor.extract_with_timeout(str(file_path))

        assert result.title == "no_title"


class TestXlsxExtractor:
    """Tests for XlsxExtractor"""

    @pytest.mark.asyncio
    async def test_extract_xlsx(self, sample_xlsx_file):
        """Test extracting from XLSX file"""
        extractor = XlsxExtractor("test_source", "Test Source")
        doc = await extractor.extract_with_timeout(sample_xlsx_file)

        assert doc.source_id == "test_source"
        assert doc.type == "xlsx"
        assert doc.basename == "sample.xlsx"
        assert doc.extension == "xlsx"
        assert "Name" in doc.content
        assert "Item One" in doc.content
        assert "100" in doc.content
        assert "Data Sheet" in doc.content
        assert "Summary" in doc.content
        assert doc.metadata.get("sheet_count") == 2
        assert doc.metadata.get("extraction_failed") is False

    @pytest.mark.asyncio
    async def test_supports_file(self, temp_dir):
        """Test file extension support"""
        assert XlsxExtractor.supports_file(str(temp_dir / "test.xlsx"))
        assert not XlsxExtractor.supports_file(str(temp_dir / "test.xls"))
        assert not XlsxExtractor.supports_file(str(temp_dir / "test.csv"))

    @pytest.mark.asyncio
    async def test_corrupted_xlsx_fallback(self, temp_dir):
        """Test fallback behavior on corrupted XLSX"""
        corrupted_xlsx = temp_dir / "corrupted.xlsx"
        corrupted_xlsx.write_text("This is not an XLSX file")

        extractor = XlsxExtractor("test_source", "Test Source")
        doc = await extractor.extract_with_timeout(str(corrupted_xlsx))

        assert doc is not None
        assert doc.content == ""
        assert doc.metadata.get("extraction_failed") is True
        assert "extraction_error" in doc.metadata

    @pytest.mark.asyncio
    async def test_empty_xlsx(self, temp_dir):
        """Test XLSX with no data"""
        from openpyxl import Workbook

        file_path = temp_dir / "empty.xlsx"
        wb = Workbook()
        wb.save(file_path)

        extractor = XlsxExtractor("test_source", "Test Source")
        doc = await extractor.extract_with_timeout(str(file_path))

        assert doc is not None
        assert doc.metadata.get("total_cells_extracted") == 0
        assert "extraction_warning" in doc.metadata


class TestPptxExtractor:
    """Tests for PptxExtractor"""

    @pytest.mark.asyncio
    async def test_extract_pptx(self, sample_pptx_file):
        """Test extracting from PPTX file"""
        extractor = PptxExtractor("test_source", "Test Source")
        doc = await extractor.extract_with_timeout(sample_pptx_file)

        assert doc.source_id == "test_source"
        assert doc.type == "pptx"
        assert doc.basename == "sample.pptx"
        assert doc.extension == "pptx"
        assert doc.title == "Test Presentation"
        assert "Presentation Title" in doc.content
        assert "Slide Two Title" in doc.content
        assert "First bullet point" in doc.content
        assert doc.metadata.get("slide_count") == 2
        assert doc.metadata.get("author") == "Test Author"
        assert doc.metadata.get("extraction_failed") is False

    @pytest.mark.asyncio
    async def test_supports_file(self, temp_dir):
        """Test file extension support"""
        assert PptxExtractor.supports_file(str(temp_dir / "test.pptx"))
        assert not PptxExtractor.supports_file(str(temp_dir / "test.ppt"))
        assert not PptxExtractor.supports_file(str(temp_dir / "test.pdf"))

    @pytest.mark.asyncio
    async def test_corrupted_pptx_fallback(self, temp_dir):
        """Test fallback behavior on corrupted PPTX"""
        corrupted_pptx = temp_dir / "corrupted.pptx"
        corrupted_pptx.write_text("This is not a PPTX file")

        extractor = PptxExtractor("test_source", "Test Source")
        doc = await extractor.extract_with_timeout(str(corrupted_pptx))

        assert doc is not None
        assert doc.content == ""
        assert doc.metadata.get("extraction_failed") is True
        assert "extraction_error" in doc.metadata

    @pytest.mark.asyncio
    async def test_pptx_without_title(self, temp_dir):
        """Test PPTX without title falls back to filename"""
        from pptx import Presentation

        file_path = temp_dir / "no_title.pptx"
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide
        prs.save(file_path)

        extractor = PptxExtractor("test_source", "Test Source")
        result = await extractor.extract_with_timeout(str(file_path))

        assert result.title == "no_title"
